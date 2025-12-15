from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Cores corporativas
COR_FUNDO = RGBColor(10, 14, 39)  # #0a0e27
COR_AZUL = RGBColor(30, 60, 114)  # #1e3c72
COR_CIANO = RGBColor(0, 212, 255)  # #00d4ff
COR_BRANCO = RGBColor(255, 255, 255)
COR_TEXTO = RGBColor(184, 197, 214)  # #b8c5d6
COR_VERDE = RGBColor(16, 185, 129)  # #10b981
COR_VERMELHO = RGBColor(239, 68, 68)  # #ef4444

def add_blank_slide(prs):
    """Adiciona slide em branco"""
    blank_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(blank_layout)
    # Fundo escuro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COR_FUNDO
    return slide

def add_header(slide, titulo, subtitulo=""):
    """Adiciona cabeçalho corporativo"""
    # Box azul do topo
    header_box = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(16), Inches(1.5)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COR_AZUL
    header_box.line.color.rgb = COR_CIANO
    header_box.line.width = Pt(4)
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COR_CIANO
    
    if subtitulo:
        p2 = tf.add_paragraph()
        p2.text = subtitulo
        p2.font.size = Pt(20)
        p2.font.color.rgb = COR_TEXTO

def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False, color=COR_BRANCO):
    """Adiciona caixa de texto"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox

def add_metric_card(slide, number, label, left, top):
    """Adiciona card de métrica"""
    # Box
    box = slide.shapes.add_shape(1, left, top, Inches(4), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    
    # Número
    add_text_box(slide, number, left, top + Inches(0.3), Inches(4), Inches(1), 60, True, COR_CIANO)
    
    # Label
    add_text_box(slide, label, left, top + Inches(1.2), Inches(4), Inches(0.7), 16, False, COR_TEXTO)

# SLIDE 1 - CAPA
slide1 = add_blank_slide(prs)
# Header azul
header_box = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(2))
header_box.fill.solid()
header_box.fill.fore_color.rgb = COR_AZUL

# Título
add_text_box(slide1, "📍 Sistema de Georreferenciamento", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), 48, True, COR_BRANCO)
add_text_box(slide1, "de Imagens para Saneamento", Inches(0.5), Inches(1.1), Inches(12), Inches(0.8), 48, True, COR_BRANCO)

# Badge prêmio
badge = slide1.shapes.add_shape(1, Inches(12.5), Inches(0.2), Inches(3), Inches(0.8))
badge.fill.solid()
badge.fill.fore_color.rgb = COR_CIANO
add_text_box(slide1, "🏆 PRÊMIO BRK\nATITUDE & INOVAÇÃO 2025", Inches(12.5), Inches(0.25), Inches(3), Inches(0.7), 14, True, COR_FUNDO)

# Subtítulo
add_text_box(slide1, "Solução tecnológica para otimização de processos operacionais\ne gestão inteligente de ativos em infraestrutura de saneamento", 
             Inches(0.5), Inches(2.3), Inches(15), Inches(1), 22, False, COR_TEXTO)

# Métricas
add_metric_card(slide1, "R$ 0", "INVESTIMENTO\nNECESSÁRIO", Inches(1), Inches(4.5))
add_metric_card(slide1, "80%", "REDUÇÃO DE TEMPO\nEM RELATÓRIOS", Inches(6), Inches(4.5))
add_metric_card(slide1, "100%", "PRECISÃO EM\nGEOLOCALIZAÇÃO", Inches(11), Inches(4.5))

# SLIDE 2 - CONTEXTO
slide2 = add_blank_slide(prs)
add_header(slide2, "Contexto Operacional no Saneamento")

# Coluna Problema
problem_box = slide2.shapes.add_shape(1, Inches(0.5), Inches(2), Inches(7), Inches(5.5))
problem_box.fill.solid()
problem_box.fill.fore_color.rgb = RGBColor(127, 29, 29)
problem_box.line.color.rgb = COR_VERMELHO
problem_box.line.width = Pt(3)

add_text_box(slide2, "❌ Desafios Atuais", Inches(0.7), Inches(2.2), Inches(6.5), Inches(0.6), 28, True, COR_BRANCO)

desafios = [
    "• Documentação Manual: Horas compilando relatórios",
    "• Perda de Informações: Fotos sem localização precisa",
    "• Retrabalho: Retornar ao campo para validar dados",
    "• Custos Elevados: Deslocamentos desnecessários",
    "• Gestão Fragmentada: Dificuldade em consolidar dados"
]
add_text_box(slide2, "\n\n".join(desafios), Inches(0.7), Inches(3), Inches(6.5), Inches(4), 16, False, COR_BRANCO)

# Coluna Solução
solution_box = slide2.shapes.add_shape(1, Inches(8.5), Inches(2), Inches(7), Inches(5.5))
solution_box.fill.solid()
solution_box.fill.fore_color.rgb = RGBColor(20, 83, 45)
solution_box.line.color.rgb = RGBColor(34, 197, 94)
solution_box.line.width = Pt(3)

add_text_box(slide2, "✅ Nossa Solução", Inches(8.7), Inches(2.2), Inches(6.5), Inches(0.6), 28, True, COR_BRANCO)

solucoes = [
    "• Automação Total: Extração automática de GPS e metadados",
    "• Rastreabilidade Completa: Coordenadas precisas",
    "• Primeira Vez Certo: Dados validados desde a captura",
    "• Otimização de Recursos: Economia em tempo e equipe",
    "• Centralização Inteligente: Plataforma única de gestão"
]
add_text_box(slide2, "\n\n".join(solucoes), Inches(8.7), Inches(3), Inches(6.5), Inches(4), 16, False, COR_BRANCO)

# SLIDE 3 - VISÃO GERAL
slide3 = add_blank_slide(prs)
add_header(slide3, "Visão Geral da Solução Tecnológica")

# Info box principal
info_box = slide3.shapes.add_shape(1, Inches(0.5), Inches(2), Inches(15), Inches(1.2))
info_box.fill.solid()
info_box.fill.fore_color.rgb = COR_AZUL
info_box.line.color.rgb = COR_CIANO
info_box.line.width = Pt(4)

add_text_box(slide3, "🔐 Plataforma Web de Alto Desempenho", Inches(0.7), Inches(2.1), Inches(14.5), Inches(0.4), 24, True, COR_CIANO)
add_text_box(slide3, "Sistema com tecnologias de ponta (Next.js 15, React 19, TypeScript) garantindo processamento local,\nsegurança de dados e conformidade LGPD. Totalmente responsivo para uso em campo.", 
             Inches(0.7), Inches(2.6), Inches(14.5), Inches(0.6), 18, False, COR_TEXTO)

# 3 boxes de características
boxes_y = Inches(4)
features = [
    ("🌐 100% Web", "Sem instalação. Acesso via\nnavegador de qualquer dispositivo"),
    ("🔒 Seguro", "Processamento local.\nDados nunca saem do dispositivo"),
    ("⚡ Rápido", "Processamento em segundos,\nmesmo para múltiplas imagens")
]

for i, (titulo, desc) in enumerate(features):
    x = Inches(0.5 + i * 5.2)
    box = slide3.shapes.add_shape(1, x, boxes_y, Inches(4.8), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    
    add_text_box(slide3, titulo, x, boxes_y + Inches(0.2), Inches(4.8), Inches(0.6), 22, True, COR_CIANO)
    add_text_box(slide3, desc, x, boxes_y + Inches(0.9), Inches(4.8), Inches(1), 16, False, COR_TEXTO)

# SLIDE 4 - FLUXO OPERACIONAL
slide4 = add_blank_slide(prs)
add_header(slide4, "Fluxo Operacional do Sistema")

steps = [
    ("1", "Captura em Campo", "Técnico fotografa com smartphone (GPS ativado)\nFormatos: JPG, PNG, HEIC"),
    ("2", "Upload na Plataforma", "Acessa via navegador e faz upload\nAté 50 fotos mobile / ilimitado desktop"),
    ("3", "Processamento Automático", "Extração instantânea de GPS, data, hora,\naltitude e thumbnail otimizado"),
    ("4", "Visualização e Edição", "Interface exibe dados. Usuário adiciona\ndescrições e ajusta status"),
    ("5", "Exportação Múltiplos Formatos", "PDF, Excel, Word, KML, JSON\nconforme necessidade")
]

y_pos = Inches(2)
for num, titulo, desc in steps:
    # Circle número
    circle = slide4.shapes.add_shape(3, Inches(0.8), y_pos, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COR_CIANO
    add_text_box(slide4, num, Inches(0.8), y_pos, Inches(0.6), Inches(0.6), 24, True, COR_FUNDO)
    
    # Content box
    content_box = slide4.shapes.add_shape(1, Inches(2), y_pos - Inches(0.1), Inches(13.5), Inches(0.8))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = COR_AZUL
    content_box.line.color.rgb = COR_CIANO
    content_box.line.width = Pt(2)
    
    add_text_box(slide4, titulo, Inches(2.2), y_pos - Inches(0.05), Inches(13), Inches(0.3), 20, True, COR_CIANO)
    add_text_box(slide4, desc, Inches(2.2), y_pos + Inches(0.28), Inches(13), Inches(0.5), 14, False, COR_TEXTO)
    
    y_pos += Inches(1.2)

# SLIDE 5 - APLICAÇÕES
slide5 = add_blank_slide(prs)
add_header(slide5, "Aplicações Práticas em Saneamento")

apps = [
    ("🔧", "Manutenção Preventiva", "Registro georeferenciado de inspeções\nem redes de água e esgoto"),
    ("💧", "Gestão de Vazamentos", "Documentação com localização exata\npara análise de padrões"),
    ("📊", "Cadastro de Ativos", "Mapeamento de hidrômetros, válvulas\ncom coordenadas precisas"),
    ("🏗️", "Acompanhamento de Obras", "Registro cronológico e geolocalizado\nde progresso de instalações"),
    ("📱", "Ordens de Serviço", "Evidências fotográficas geolocalizadas\nde serviços executados"),
    ("🗺️", "Integração SIG", "Exportação KML para Google Earth\ne sistemas GIS corporativos")
]

for i, (icon, titulo, desc) in enumerate(apps):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 5.2)
    y = Inches(2 + row * 2.5)
    
    box = slide5.shapes.add_shape(1, x, y, Inches(4.8), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    
    add_text_box(slide5, icon, x + Inches(2), y + Inches(0.1), Inches(1), Inches(0.5), 36, False, COR_BRANCO)
    add_text_box(slide5, titulo, x + Inches(0.2), y + Inches(0.7), Inches(4.4), Inches(0.4), 18, True, COR_CIANO)
    add_text_box(slide5, desc, x + Inches(0.2), y + Inches(1.2), Inches(4.4), Inches(0.7), 14, False, COR_TEXTO)

# SLIDE 6 - BENEFÍCIOS E ROI
slide6 = add_blank_slide(prs)
add_header(slide6, "Benefícios Mensuráveis e ROI")

# Ganhos operacionais (esquerda)
add_text_box(slide6, "Ganhos Operacionais", Inches(0.5), Inches(2), Inches(7), Inches(0.5), 24, True, COR_CIANO)

ganhos = [
    "⏱️  80% de redução no tempo de elaboração de relatórios",
    "🎯  100% de precisão na localização de ativos",
    "🚗  Redução de deslocamentos desnecessários",
    "📈  Aumento de produtividade das equipes",
    "🌱  Sustentabilidade: menos impressões e combustível"
]

y_ganhos = Inches(2.7)
for ganho in ganhos:
    box = slide6.shapes.add_shape(1, Inches(0.5), y_ganhos, Inches(7), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    add_text_box(slide6, ganho, Inches(0.7), y_ganhos + Inches(0.15), Inches(6.5), Inches(0.5), 16, False, COR_BRANCO)
    y_ganhos += Inches(0.9)

# ROI Box (direita)
roi_box = slide6.shapes.add_shape(1, Inches(8.5), Inches(2), Inches(7), Inches(3))
roi_box.fill.solid()
roi_box.fill.fore_color.rgb = RGBColor(6, 95, 70)
roi_box.line.color.rgb = COR_VERDE
roi_box.line.width = Pt(4)

add_text_box(slide6, "Investimento Necessário", Inches(8.5), Inches(2.2), Inches(7), Inches(0.5), 28, True, COR_BRANCO)
add_text_box(slide6, "R$ 0,00", Inches(8.5), Inches(2.9), Inches(7), Inches(0.9), 72, True, COR_VERDE)

itens_roi = "✅ Licença gratuita\n✅ Sem limite de usuários\n✅ Sem limite de processamento\n✅ Atualizações incluídas\n✅ Suporte da comunidade\n✅ Hospedagem inclusa"
add_text_box(slide6, itens_roi, Inches(8.7), Inches(4), Inches(6.5), Inches(1), 16, False, RGBColor(209, 250, 229))

# Análise comparativa
comp_box = slide6.shapes.add_shape(1, Inches(8.5), Inches(5.5), Inches(7), Inches(1.5))
comp_box.fill.solid()
comp_box.fill.fore_color.rgb = COR_AZUL
comp_box.line.color.rgb = COR_CIANO
comp_box.line.width = Pt(2)

add_text_box(slide6, "💰 Análise Comparativa", Inches(8.7), Inches(5.65), Inches(6.5), Inches(0.4), 20, True, COR_CIANO)
add_text_box(slide6, "Soluções comerciais: R$ 500 a R$ 2.000/usuário/ano\nEquipe de 50 usuários: economia de até R$ 100.000/ano", 
             Inches(8.7), Inches(6.1), Inches(6.5), Inches(0.8), 14, False, COR_TEXTO)

# SLIDE 7 - TECNOLOGIA E SEGURANÇA
slide7 = add_blank_slide(prs)
add_header(slide7, "Arquitetura Tecnológica e Segurança")

# Stack Tecnológico (esquerda)
add_text_box(slide7, "Stack Tecnológico", Inches(0.5), Inches(2), Inches(7), Inches(0.5), 24, True, COR_CIANO)

stack = [
    "⚛️  Next.js 15 + React 19: Framework de alta performance",
    "📘  TypeScript: Código robusto e tipagem estática",
    "🎨  Tailwind CSS: Interface responsiva e profissional",
    "📦  ExifReader: Extração eficiente de metadados EXIF",
    "☁️  Vercel Edge Network: Deploy global com baixa latência"
]

y_stack = Inches(2.7)
for item in stack:
    box = slide7.shapes.add_shape(1, Inches(0.5), y_stack, Inches(7), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    add_text_box(slide7, item, Inches(0.7), y_stack + Inches(0.15), Inches(6.5), Inches(0.5), 15, False, COR_BRANCO)
    y_stack += Inches(0.9)

# Segurança (direita)
add_text_box(slide7, "Segurança e Conformidade", Inches(8.5), Inches(2), Inches(7), Inches(0.5), 24, True, COR_CIANO)

seg_boxes = [
    ("🔐 Privacidade Total", "Processamento 100% no navegador. Imagens nunca transmitidas.\nConformidade total com LGPD."),
    ("🛡️ Segurança de Dados", "HTTPS obrigatório. Sem armazenamento em nuvem.\nUsuário mantém controle total."),
    ("✅ Disponibilidade", "SLA 99.9% garantido. CDN global.\nFuncionamento offline após carregamento.")
]

y_seg = Inches(2.7)
for titulo, desc in seg_boxes:
    box = slide7.shapes.add_shape(1, Inches(8.5), y_seg, Inches(7), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    add_text_box(slide7, titulo, Inches(8.7), y_seg + Inches(0.1), Inches(6.5), Inches(0.4), 18, True, COR_CIANO)
    add_text_box(slide7, desc, Inches(8.7), y_seg + Inches(0.5), Inches(6.5), Inches(0.7), 14, False, COR_TEXTO)
    y_seg += Inches(1.5)

# SLIDE 8 - DIFERENCIAIS
slide8 = add_blank_slide(prs)
add_header(slide8, "Diferenciais Competitivos")

# 3 cards principais
difs_principais = [
    ("🚀", "Velocidade", "Processamento instantâneo vs.\nhoras de trabalho manual"),
    ("💎", "Qualidade", "Dados estruturados e\npadronizados automaticamente"),
    ("🌍", "Acessibilidade", "Uso em qualquer lugar,\nqualquer dispositivo")
]

y_top = Inches(2)
for i, (icon, titulo, desc) in enumerate(difs_principais):
    x = Inches(0.5 + i * 5.2)
    box = slide8.shapes.add_shape(1, x, y_top, Inches(4.8), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    
    add_text_box(slide8, icon, x + Inches(2), y_top + Inches(0.1), Inches(1), Inches(0.5), 48, False, COR_BRANCO)
    add_text_box(slide8, titulo, x + Inches(0.2), y_top + Inches(0.8), Inches(4.4), Inches(0.4), 22, True, COR_BRANCO)
    add_text_box(slide8, desc, x + Inches(0.2), y_top + Inches(1.3), Inches(4.4), Inches(0.6), 16, False, COR_TEXTO)

# Lista de diferenciais
add_text_box(slide8, "Por Que Escolher Esta Solução?", Inches(0.5), Inches(4.5), Inches(15), Inches(0.5), 26, True, COR_CIANO)

difs_lista = [
    ("💰 Custo Zero", "🎯 Fácil Adoção"),
    ("📱 Mobile First", "🔌 Interoperabilidade"),
    ("🔄 Evolução Contínua", "🌐 Escalável")
]

y_lista = Inches(5.2)
for esq, dir in difs_lista:
    # Esquerda
    box_e = slide8.shapes.add_shape(1, Inches(0.5), y_lista, Inches(7.2), Inches(0.7))
    box_e.fill.solid()
    box_e.fill.fore_color.rgb = COR_AZUL
    box_e.line.color.rgb = COR_CIANO
    box_e.line.width = Pt(2)
    add_text_box(slide8, esq, Inches(0.7), y_lista + Inches(0.15), Inches(6.8), Inches(0.5), 18, False, COR_BRANCO)
    
    # Direita
    box_d = slide8.shapes.add_shape(1, Inches(8.3), y_lista, Inches(7.2), Inches(0.7))
    box_d.fill.solid()
    box_d.fill.fore_color.rgb = COR_AZUL
    box_d.line.color.rgb = COR_CIANO
    box_d.line.width = Pt(2)
    add_text_box(slide8, dir, Inches(8.5), y_lista + Inches(0.15), Inches(6.8), Inches(0.5), 18, False, COR_BRANCO)
    
    y_lista += Inches(0.9)

# SLIDE 9 - ROADMAP
slide9 = add_blank_slide(prs)
add_header(slide9, "Roadmap de Evolução 2025-2026")

roadmap = [
    ("Q1", "Mapa Interativo", "Visualização de imagens em mapa com clustering e filtros"),
    ("Q2", "Medições Geodésicas", "Cálculo de distâncias, áreas e perímetros"),
    ("Q3", "IA para Análise", "Detecção automática de anomalias em infraestrutura"),
    ("Q4", "API Corporativa", "Integração com sistemas ERP/SAP e GIS corporativos")
]

y_road = Inches(2.2)
for trimestre, titulo, desc in roadmap:
    # Circle trimestre
    circle = slide9.shapes.add_shape(3, Inches(0.8), y_road, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COR_CIANO
    add_text_box(slide9, trimestre, Inches(0.8), y_road + Inches(0.05), Inches(0.7), Inches(0.6), 20, True, COR_FUNDO)
    
    # Content
    box = slide9.shapes.add_shape(1, Inches(2), y_road, Inches(6), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    add_text_box(slide9, titulo, Inches(2.2), y_road + Inches(0.1), Inches(5.5), Inches(0.4), 20, True, COR_CIANO)
    add_text_box(slide9, desc, Inches(2.2), y_road + Inches(0.5), Inches(5.5), Inches(0.4), 15, False, COR_TEXTO)
    
    y_road += Inches(1.3)

# Visão de longo prazo (direita)
visao_box = slide9.shapes.add_shape(1, Inches(8.5), Inches(2.2), Inches(7), Inches(5.5))
visao_box.fill.solid()
visao_box.fill.fore_color.rgb = COR_AZUL
visao_box.line.color.rgb = COR_CIANO
visao_box.line.width = Pt(2)

add_text_box(slide9, "🎯 Visão de Longo Prazo", Inches(8.7), Inches(2.4), Inches(6.5), Inches(0.5), 22, True, COR_CIANO)
add_text_box(slide9, "Tornar-se a principal plataforma open-source de\ngeorreferenciamento para o setor de saneamento no Brasil:", 
             Inches(8.7), Inches(3), Inches(6.5), Inches(0.8), 16, False, COR_TEXTO)

visao_itens = [
    "🤝  Comunidade ativa de desenvolvedores",
    "📚  Treinamentos e certificações",
    "🌎  Expansão para outros segmentos",
    "🔬  Parcerias com universidades"
]

y_visao = Inches(4)
for item in visao_itens:
    add_text_box(slide9, item, Inches(8.7), y_visao, Inches(6.5), Inches(0.5), 16, False, COR_BRANCO)
    y_visao += Inches(0.7)

# SLIDE 10 - CONCLUSÃO
slide10 = add_blank_slide(prs)
add_header(slide10, "Inovação Acessível que Gera Valor Real")

# 3 métricas principais
add_metric_card(slide10, "R$ 0", "INVESTIMENTO", Inches(1), Inches(2.5))
add_metric_card(slide10, "∞", "POTENCIAL\nDE ROI", Inches(6), Inches(2.5))
add_metric_card(slide10, "100%", "ATITUDE\nBRK", Inches(11), Inches(2.5))

# Box destaque
destaque_box = slide10.shapes.add_shape(1, Inches(1.5), Inches(5), Inches(13), Inches(2))
destaque_box.fill.solid()
destaque_box.fill.fore_color.rgb = RGBColor(6, 95, 70)
destaque_box.line.color.rgb = COR_VERDE
destaque_box.line.width = Pt(4)

add_text_box(slide10, "Esta é a Atitude da Inovação", Inches(1.5), Inches(5.2), Inches(13), Inches(0.6), 36, True, COR_BRANCO)
add_text_box(slide10, "Democratizar tecnologia de ponta, eliminar barreiras de custo,\notimizar processos operacionais e gerar eficiência mensurável\npara o setor de saneamento brasileiro.", 
             Inches(1.5), Inches(6), Inches(13), Inches(1), 20, False, RGBColor(209, 250, 229))

# Link de acesso
add_text_box(slide10, "Acesse e Teste Agora", Inches(0.5), Inches(7.5), Inches(15), Inches(0.4), 28, True, COR_CIANO)

link_box = slide10.shapes.add_shape(1, Inches(4.5), Inches(8), Inches(7), Inches(0.6))
link_box.fill.solid()
link_box.fill.fore_color.rgb = COR_AZUL
link_box.line.color.rgb = COR_CIANO
link_box.line.width = Pt(3)
add_text_box(slide10, "🔗 relatoriopendv0.vercel.app", Inches(4.5), Inches(8.05), Inches(7), Inches(0.5), 24, True, COR_CIANO)

# SLIDE 11 - AGRADECIMENTO
slide11 = add_blank_slide(prs)

add_text_box(slide11, "Obrigado!", Inches(0.5), Inches(2.5), Inches(15), Inches(1), 72, True, COR_BRANCO)
add_text_box(slide11, "Comprometidos com a excelência operacional\ne a transformação digital do saneamento brasileiro", 
             Inches(0.5), Inches(3.8), Inches(15), Inches(1), 26, False, COR_TEXTO)

# 3 cards finais
final_cards = [
    ("📧", "Contato", "Dúvidas técnicas e\noperacionais"),
    ("💡", "Sugestões", "Sua experiência\nnos melhora"),
    ("🏆", "BRK Atitude", "Inovação que\ntransforma")
]

for i, (icon, titulo, desc) in enumerate(final_cards):
    x = Inches(1.5 + i * 4.5)
    y = Inches(5.5)
    
    box = slide11.shapes.add_shape(1, x, y, Inches(4), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_AZUL
    box.line.color.rgb = COR_CIANO
    box.line.width = Pt(2)
    
    add_text_box(slide11, icon, x + Inches(1.5), y + Inches(0.1), Inches(1), Inches(0.5), 48, False, COR_BRANCO)
    add_text_box(slide11, titulo, x + Inches(0.2), y + Inches(0.8), Inches(3.6), Inches(0.4), 20, True, COR_CIANO)
    add_text_box(slide11, desc, x + Inches(0.2), y + Inches(1.3), Inches(3.6), Inches(0.6), 16, False, COR_TEXTO)

# Salvar apresentação
output_path = os.path.join(os.path.dirname(__file__), 'Apresentacao_BRK_Atitude_Inovacao_2025.pptx')
prs.save(output_path)
print(f"✅ Apresentação criada com sucesso!")
print(f"📁 Arquivo: {output_path}")
